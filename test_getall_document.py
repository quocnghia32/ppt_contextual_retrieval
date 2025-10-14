from pptx import Presentation

def iter_text_in_shape(shape):
    if not hasattr(shape, "has_text_frame") or not shape.has_text_frame:
        return []
    lines = []
    for p in shape.text_frame.paragraphs:
        level = p.level or 0
        text = ''.join(run.text for run in p.runs).strip()
        if text:
            lines.append(("  " * level) + f"- {text}")
    return lines

def extract_table(shape):
    tbl = shape.table
    rows = []
    for r in tbl.rows:
        cells = [c.text.strip().replace("\n", " ") for c in r.cells]
        rows.append(" | ".join(cells))
    return ["[TABLE]"] + rows

def slide_title(slide):
    # title placeholder if present
    if slide.shapes.title and slide.shapes.title.text:
        return slide.shapes.title.text.strip()
    # fallback: first non-empty text shape
    for shp in slide.shapes:
        if getattr(shp, "has_text_frame", False) and shp.text_frame and shp.text_frame.text.strip():
            return shp.text_frame.text.strip().split("\n", 1)[0]
    return "Untitled"

def extract_notes(slide):
    if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
        t = slide.notes_slide.notes_text_frame.text.strip()
        return ["[NOTES]", t] if t else []
    return []

def whole_document_from_pptx(path, deck_title=None):
    prs = Presentation(path)
    deck_title = deck_title or (prs.core_properties.title or "")
    parts = []
    parts.append(f"[DECK_TITLE] {deck_title or path}")
    for idx, slide in enumerate(prs.slides, start=1):
        parts.append(f"\nSlide {idx} — {slide_title(slide)}")
        # visible text & tables
        for shp in slide.shapes:
            if getattr(shp, "has_text_frame", False) and shp.has_text_frame:
                parts.extend(iter_text_in_shape(shp))
            elif hasattr(shp, "has_table") and shp.has_table:
                parts.extend(extract_table(shp))
            # (Optional) alt text
            alt = getattr(shp, "alternative_text", "") or ""
            if alt.strip():
                parts.append(f"[ALT] {alt.strip()}")
        # speaker notes
        parts.extend(extract_notes(slide))
    return "\n".join(parts)

if __name__ == "__main__":
    print(whole_document_from_pptx("test.pptx"))
