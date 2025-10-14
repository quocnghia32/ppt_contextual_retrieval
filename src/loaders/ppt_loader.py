"""
Custom LangChain document loader for PowerPoint files.
"""
import os
from langchain.document_loaders.base import BaseLoader
from langchain.schema import Document
from typing import List, Optional, Tuple
import hashlib
import json
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import io
from loguru import logger

from src.models.vision_analyzer import VisionAnalyzer
from src.utils.image_checker import is_valid_image
class PPTLoader(BaseLoader):
    """
    LangChain Document Loader for PowerPoint (.pptx) files.

    Extracts slides with metadata and optionally extracts images.
    """

    def __init__(
        self,
        file_path: str,
        extract_images: bool = True,
        include_speaker_notes: bool = True,
        use_vision: bool = True,
        extract_charts: bool = True,
        extract_tables: bool = True
    ):
        """
        Initialize PPT Loader.

        Args:
            file_path: Path to .pptx file
            extract_images: Whether to extract images from slides
            include_speaker_notes: Whether to include speaker notes
            extract_charts: Whether to extract chart data
            extract_tables: Whether to extract table data
        """
        self.file_path = file_path
        self.extract_images = extract_images
        self.include_speaker_notes = include_speaker_notes
        self.extract_charts = extract_charts
        self.extract_tables = extract_tables
        self.vision_analyzer = VisionAnalyzer() if use_vision else None

    async def load(self) -> Tuple[List[Document], Document]:
        """
        Load PPT and return list of LangChain Documents.

        Each slide becomes a Document with comprehensive metadata.
        """
        try:
            prs = Presentation(self.file_path)
            documents = []

            # Extract presentation metadata
            presentation_id = hashlib.md5(self.file_path.encode()).hexdigest()[:12]
            title = self._extract_presentation_title(prs)

            logger.info(f"Loading presentation: {title} ({len(prs.slides)} slides)")

            # Detect sections (if any)
            sections = self._detect_sections(prs)

            logger.info(f"Detected {len(sections)} sections")
           
            # Add a overall information document about current file
            overall_info = Document(
                page_content=self._extract_ppt_info(prs),
                metadata={"source": self.file_path,
                          "slide_number": 0,
                          
                          }
            )

            for slide_idx, slide in enumerate(prs.slides):
                slide_num = slide_idx + 1

                # Extract slide content
                slide_text = self._extract_slide_text(slide)
                slide_title = self._extract_slide_title(slide)

                # Extract speaker notes
                speaker_notes = ""
                if self.include_speaker_notes and slide.has_notes_slide:
                    try:
                        speaker_notes = slide.notes_slide.notes_text_frame.text
                    except Exception as e:
                        logger.warning(f"Failed to extract notes from slide {slide_num}: {e}")

                # Extract images
                images_info = []
                if self.extract_images:
                    images_info = await self._extract_images_info(slide, slide_num)


                # Extract charts
                charts_info = []
                if self.extract_images:
                    charts_info = self._extract_charts_info(slide, slide_num)

                # Extract tables
                tables_info = []
                if self.extract_tables:
                    tables_info = self._extract_tables(slide)

                # Determine section
                section = self._get_section_for_slide(slide_num, sections, prs)

                # Build comprehensive metadata
                # Convert complex objects to JSON strings for Pinecone compatibility
                metadata = {
                    "source": self.file_path,
                    "presentation_id": presentation_id,
                    "presentation_title": title,
                    "slide_number": slide_num,
                    "total_slides": len(prs.slides),
                    "slide_title": slide_title,
                    "section": section,
                    "speaker_notes": speaker_notes,
                    "image_count": len(images_info),
                    "images": json.dumps(images_info),  # Convert to JSON string
                    "chart_count": len(charts_info),
                    "charts": json.dumps(charts_info),  # Convert to JSON string
                    "table_count": len(tables_info),
                    "tables": json.dumps(tables_info),  # Convert to JSON string
                    "type": "slide"
                }

                # Combine all text content
                full_content = self._build_full_content(
                    slide_text,
                    speaker_notes,
                    tables_info,
                    images_info
                )

                # Create Document
                doc = Document(
                    page_content=full_content,
                    metadata=metadata
                )
                documents.append(doc)

            logger.info(f"Successfully loaded {len(documents)} slides")
            return documents, overall_info

        except Exception as e:
            logger.error(f"Failed to load PPT {self.file_path}: {e}")
            raise

    def _extract_ppt_info(self, prs: Presentation) -> str:
        """
        Extracts simple overall information from a PPTX presentation.
        Includes file name, total slides, section names, and slide titles.
        """
        props = prs.core_properties
        total_slides = len(prs.slides)
        file_name = os.path.basename(self.file_path)

        title = props.title or file_name
        author = props.author or "Unknown author"
        created = props.created.strftime("%Y-%m-%d") if props.created else "Unknown"
        modified = props.modified.strftime("%Y-%m-%d") if props.modified else "Unknown"

        # --- Try to get section names ---
        # python-pptx doesn’t expose sections directly, but you can read them via the underlying XML
        section_names = []
        try:
            sldIdLst = prs.part.element.xpath("//p:sldIdLst/p:sldId")
            sections = prs.part.element.xpath("//p:sectionLst/p:section")
            for s in sections:
                name = s.get("name")
                if name:
                    section_names.append(name)
        except Exception:
            section_names = []

        # --- Slide titles ---
        slide_titles = []
        for i, slide in enumerate(prs.slides, start=1):
            title_shape = slide.shapes.title
            if title_shape and title_shape.text.strip():
                title_text = title_shape.text.strip()
            else:
                title_text = "(No title)"
            slide_titles.append(f"{i}. {title_text}")

        # --- Build summary text ---
        info_text = (
            f"OVERALL INFORMATION ABOUT THE PRESENTATION/FILE/SLIDES:\n"
            f"File name: {file_name}\n"
            f"Presentation title: {title}\n"
            f"Author: {author}\n"
            f"Total slides: {total_slides}\n"
            f"Created: {created}\n"
            f"Modified: {modified}\n"
        )

        if section_names:
            info_text += "Sections:\n" + "\n".join(f"- {s}" for s in section_names) + "\n"

        info_text += "\nSlide titles:\n" + "\n".join(slide_titles)

        return info_text


    def _extract_presentation_title(self, prs: Presentation) -> str:
        """Extract presentation title from first slide."""
        if len(prs.slides) > 0:
            first_slide = prs.slides[0]
            if first_slide.shapes.title:
                return first_slide.shapes.title.text.strip()
        return "Untitled Presentation"

    def _extract_slide_title(self, slide) -> str:
        """Extract title from slide."""
        try:
            if slide.shapes.title:
                return slide.shapes.title.text.strip()
        except Exception:
            pass
        return ""

    def _extract_slide_text(self, slide) -> str:
        """Extract all text content from slide."""
        text_parts = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_parts.append(shape.text.strip())

        return "\n\n".join(text_parts)
    
    def _extract_charts_info(self, slide, slide_num: int) -> List[dict]:
        """Extract chart data from slide."""
        charts = []

        for shape_idx, shape in enumerate(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                charts.append({
                    "slide_number": slide_num,
                    "shape_index": shape_idx,
                    "type": "chart",
                    "shape_id": shape.shape_id
                })

        return charts

    async def _extract_images_info(self, slide, slide_num: int) -> List[dict]:
        """
        Extract information about images in slide.
        """
        images = []
        # data = f"""
        #             **Context Information:**
        #             - Slide Number: {slide_context.get('slide_number')}
        #             - Section: {slide_context.get('section', 'Unknown')}
        #             - Slide Title: {slide_context.get('slide_title', 'Unknown')}
        #             - Total Slides: {slide_context.get('total_slides', 'Unknown')}
        #             """
        success_count = 0

        for shape_idx, shape in enumerate(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    image_bytes = image.blob
                    image_ext = image.ext
                    if not is_valid_image(image_bytes):
                        logger.warning(f"Skipping invalid image with shape index {shape_idx} in slide {slide_num}")
                        continue

                    # Get image info
                    image_info = {
                        "slide_number": slide_num,
                        "image_index": shape_idx,
                        "type": "picture",
                        # Store reference for later processing
                        "shape_id": shape.shape_id,
                        "image_format": image_ext,
                        "image_size": len(image_bytes)
                    }
                    vision_description = await self.vision_analyzer.process_image(self.file_path, image_info, image_bytes)
                    if vision_description is not None:
                        # Add image metadata
                        success_count += 1
                        image_info['vision_description'] = vision_description

                    # Try to get dimensions
                    if hasattr(shape, "width") and hasattr(shape, "height"):
                        image_info["width"] = shape.width
                        image_info["height"] = shape.height

                    images.append(image_info)

                except Exception as e:
                    logger.warning(f"Failed to process image in slide {slide_num}: {e}")
        logger.debug(f"Analyzed {success_count} images in slide {slide_num}")
        return images

    def _extract_tables(self, slide) -> List[dict]:
        """Extract table data from slide."""
        tables = []

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                try:
                    table = shape.table
                    rows = []

                    for row in table.rows:
                        row_data = []
                        for cell in row.cells:
                            row_data.append(cell.text.strip())
                        rows.append(row_data)

                    tables.append({
                        "rows": rows,
                        "row_count": len(table.rows),
                        "col_count": len(table.columns)
                    })

                except Exception as e:
                    logger.warning(f"Failed to extract table: {e}")

        return tables

    def _detect_sections(self, prs: Presentation) -> List[dict]:
        """
        Detect section boundaries based on slide titles.

        Simple heuristic: Look for slides that look like section headers.
        """
        sections = []

        section_keywords = [
            "introduction", "overview", "agenda",
            "background", "conclusion", "summary",
            "results", "analysis", "recommendations"
        ]

        for idx, slide in enumerate(prs.slides):
            title = self._extract_slide_title(slide).lower()

            # Check if this looks like a section header
            is_section_header = False

            # Heuristic 1: Contains section keywords
            if any(keyword in title for keyword in section_keywords):
                is_section_header = True

            # Heuristic 2: Slide has only title (section divider)
            text = self._extract_slide_text(slide)
            if title and len(text.split()) < 20:  # Very short content
                is_section_header = True

            if is_section_header:
                sections.append({
                    "title": self._extract_slide_title(slide),
                    "start_slide": idx + 1
                })

        return sections

    def _get_section_for_slide(
        self,
        slide_num: int,
        sections: List[dict],
        prs: Presentation
    ) -> str:
        """Determine which section a slide belongs to."""
        if not sections:
            # Default sections
            total_slides = len(prs.slides)
            if slide_num == 1:
                return "Introduction"
            elif slide_num == total_slides:
                return "Conclusion"
            else:
                return "Main Content"

        # Find the section this slide belongs to
        current_section = "Introduction"
        for section in sections:
            if slide_num >= section["start_slide"]:
                current_section = section["title"]
            else:
                break

        return current_section

    def _build_full_content(
        self,
        slide_text: str,
        speaker_notes: str,
        tables: List[dict],
        images: List[dict]
    ) -> str:
        """Build full content string from all text sources."""
        parts = []

        # Add main slide text
        if slide_text:
            parts.append(slide_text)

        # Add table data
        for table in tables:
            table_text = self._format_table(table["rows"])
            parts.append(f"\n[Table]\n{table_text}")

        # Add speaker notes
        if speaker_notes:
            parts.append(f"\n[Speaker Notes]\n{speaker_notes}")

        
        image_texts = []
        for idx, img in enumerate(images, 1):
            img_text_parts = [f"\n\n[IMAGE {idx}]"]

            # Add vision description
            if img.get("vision_description"):
                img_text_parts.append(f"Description: {img['vision_description']}")

            image_texts.append(" ".join(img_text_parts))

        # Append all image info to page content
        parts.append("".join(image_texts))

        return "\n\n".join(parts)

    def _format_table(self, rows: List[List[str]]) -> str:
        """Format table as markdown."""
        if not rows:
            return ""

        # Simple markdown table
        lines = []
        for row in rows:
            lines.append("| " + " | ".join(row) + " |")

        # Add header separator after first row
        if len(lines) > 1:
            col_count = len(rows[0])
            separator = "| " + " | ".join(["---"] * col_count) + " |"
            lines.insert(1, separator)

        return "\n".join(lines)
