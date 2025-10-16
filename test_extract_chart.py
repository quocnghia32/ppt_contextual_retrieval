from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def get_slide_shape(container):
    """
    Yield all shapes in a slide or group recursively.
    Works for slide-level containers and group shapes.
    """
    for shape in container.shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            # recurse into child shapes of the group
            yield from get_slide_shape(shape)
def get_all_text(shapes):
    text_parts = []

    for shape in shapes:
        if hasattr(shape, "text") and shape.text.strip():
            text_parts.append(shape.text.strip())

    return "\n1\n".join(text_parts)
def extract_ppt_type(ppt_path):
    mapne = {}
    prs = Presentation(ppt_path)
    for i, slide in enumerate(prs.slides,1):
        
        all_shape = list(get_slide_shape(slide))
        for shape in all_shape:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                if i==6:
                    print("*"*50)
            if hasattr(shape, "text") and shape.text.strip():
                if i==6:
                    print(shape.text.strip())
            if i==6:
                print(shape.shape_type)
                mapne[shape.shape_type] = mapne.get(shape.shape_type, 0) + 1
        if i==6:
            print("======")
            print(get_all_text(all_shape))
            print("======")

    for k,v in mapne.items():
        print(k,v)


if __name__ == "__main__":
    extract_ppt_type("hihi.pptx")